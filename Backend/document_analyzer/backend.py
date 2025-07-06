from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
import os
import random
import string
import json
import tempfile
import shutil

from docx import Document
from pypdf import PdfReader
from pptx import Presentation

from quiz_generator import export_quiz
from flash_card_generator import export_flashcards
from summary_generator import export_summary
from gemini_call import prompt_everyting
import google.generativeai as genai
from goal_todo_features import GoalTodoManager
from video_search import VideoSearchManager

# Ensure required directories exist
os.makedirs("database", exist_ok=True)

app = Flask(__name__)

# --- ✅ Fix for CORS ---
CORS(app, origins=["http://localhost:5173","http://localhost:8080"], supports_credentials=True)

# --- Gemini Chatbot Setup ---
class Chatbot:
    def __init__(self):
        self.api_key = "AIzaSyDSRv_0xjtnd92cCnvuCAv7QB-1PJcVU1Y"
        genai.configure(api_key=self.api_key)
        self.model = genai.GenerativeModel("gemini-2.0-flash")
        self.chat = None
        self.content = None

    def start_chat(self, content: str) -> str:
        try:
            prompt = f"You are a helpful assistant that answers questions based only on the provided content:\n{content}"
            self.chat = self.model.start_chat(history=[])
            self.chat.send_message(prompt)
            self.content = content
            return "Chat started successfully."
        except Exception as e:
            return f"Error starting chat: {str(e)}"

    def get_response(self, message: str) -> str:
        try:
            if not self.chat:
                return "Please start a chat with content first."
            response = self.chat.send_message(message)
            return response.text
        except Exception as e:
            return f"Error generating response: {str(e)}"

    def clear_chat(self) -> str:
        self.chat = None
        self.content = None
        return "Chat cleared."

chatbot = Chatbot()
goal_todo_manager = GoalTodoManager(api_key=chatbot.api_key)
video_search_manager = VideoSearchManager()

# --- Utility Functions ---
def generate_quiz(flashcards):
    quiz = []
    for card in flashcards:
        question = {"question": card[0], "possible_answers": [], "index": -1}
        incorrect = [c[1] for c in flashcards if c != card]
        question["possible_answers"] = random.sample(incorrect, min(3, len(incorrect)))
        question["index"] = random.randint(0, len(question["possible_answers"]))
        question["possible_answers"].insert(question["index"], card[1])
        quiz.append(question)
    return quiz

def generate_id():
    existing = {f[:-5] for f in os.listdir("./database") if f.endswith(".json")}
    while True:
        new_id = ''.join(random.choices(string.ascii_letters + string.digits, k=5))
        if new_id not in existing:
            return new_id

def handle_pdf(file_path): return ''.join(PdfReader(file_path).pages[i].extract_text() for i in range(len(PdfReader(file_path).pages)))
def handle_txt(file_path): return open(file_path, 'r', encoding='utf-8').read()
def handle_docx(file_path): return '\n'.join(p.text for p in Document(file_path).paragraphs)
def handle_pptx(file_path): return '\n'.join(shape.text for slide in Presentation(file_path).slides for shape in slide.shapes if hasattr(shape, "text"))

# --- Routes ---
@app.route("/")
def hello(): return "Hello, World!"

@app.route("/upload", methods=["POST"])
def upload():
    temp_dir = None
    try:
        file = request.files.get("file")
        if not file:
            return jsonify({"error": "No file provided"}), 400

        ext = file.filename.rsplit(".", 1)[-1].lower()
        temp_dir = tempfile.mkdtemp()
        file_path = os.path.join(temp_dir, f"file.{ext}")
        file.save(file_path)

        if ext == "pdf":
            content = handle_pdf(file_path)
        elif ext == "txt":
            content = handle_txt(file_path)
        elif ext == "docx":
            content = handle_docx(file_path)
        elif ext == "pptx":
            content = handle_pptx(file_path)
        else:
            return jsonify({"error": "Unsupported file format"}), 400

        if not content.strip():
            return jsonify({"error": "No text extracted from file"}), 400

        response = prompt_everyting(content)
        if not response:
            return jsonify({"error": "AI failed to process content"}), 500

        response["quiz"] += generate_quiz(response["flash_cards"])
        random.shuffle(response["quiz"])
        response["id"] = generate_id()
        chatbot.start_chat(content)

        db_path = os.path.join("database", f"{response['id']}.json")
        with open(db_path, "w", encoding="utf-8") as f:
            json.dump(response, f, ensure_ascii=False, indent=2)

        return jsonify(response), 200

    except Exception as e:
        print(f"Upload error: {e}")
        return jsonify({"error": f"Upload failed: {str(e)}"}), 500
    finally:
        if temp_dir and os.path.exists(temp_dir):
            shutil.rmtree(temp_dir, ignore_errors=True)

@app.route("/fetch_id", methods=["POST"])
def fetch_id():
    id = request.json.get("id")
    try:
        with open(f"database/{id}.json", "r", encoding="utf-8") as f:
            return json.load(f)
    except:
        return {"summary": 404}

@app.route("/recent", methods=["GET"])
def recent():
    files = sorted([f for f in os.listdir("database") if f.endswith(".json")], reverse=True)[:10]
    output = {"recent": []}
    for file in files:
        try:
            with open(f"database/{file}", "r", encoding="utf-8") as f:
                data = json.load(f)
                output["recent"].append({"id": file[:-5], "title": data.get("title", "Untitled")})
        except Exception as e:
            print(f"Error reading {file}: {e}")
    return output

@app.route("/export", methods=["POST"])
def export():
    selected = request.json.get("selected")
    data = request.json.get("data")
    if selected == 0:
        export_summary(data["summary"], "Summary.docx")
        return send_file("Summary.docx", as_attachment=True)
    elif selected == 1:
        export_flashcards(data["flash_cards"], "Flashcards.docx")
        return send_file("Flashcards.docx", as_attachment=True)
    else:
        export_quiz(data["quiz"], "Quiz.docx")
        return send_file("Quiz.docx", as_attachment=True)

@app.route("/chat/start", methods=["POST"])
def start_chat():
    content = request.json.get("content")
    if not content:
        return jsonify({"error": "No content provided"}), 400
    result = chatbot.start_chat(content)
    return jsonify({"message": result})

@app.route("/chat/message", methods=["POST"])
def chat_message():
    message = request.json.get("message")
    if not message:
        return jsonify({"error": "No message provided"}), 400
    return jsonify({"response": chatbot.get_response(message)})

@app.route("/chat/clear", methods=["POST"])
def clear_chat():
    return jsonify({"message": chatbot.clear_chat()})

@app.route("/goaladvise", methods=["POST"])
def goal_advise():
    goal = request.json.get("goal")
    if not goal:
        return jsonify({"error": "No goal statement provided"}), 400
    try:
        return jsonify(goal_todo_manager.get_goal_advice(goal)), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/generatetodotask", methods=["POST"])
def generate_todo_auto():
    goal = request.json.get("goal")
    if not goal:
        return jsonify({"error": "No goal provided"}), 400
    try:
        return jsonify(goal_todo_manager.generate_todo_tasks(goal)), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/search_videos", methods=["POST"])
def search_videos():
    query = request.json.get("query")
    if not query:
        return jsonify({"error": "No search query provided"}), 400
    try:
        return jsonify(video_search_manager.format_video_response(query)), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    app.run(debug=True, port=5002)
